"""
File parser for goal discovery.

Extracts text content from various file formats so it can be sent to the LLM.
Supports: md, txt, json, csv, docx, pptx, xlsx, jpg/jpeg/png (as images), zip.
"""

import base64
import io
import json
import os
import zipfile
from dataclasses import dataclass, field
from typing import List, Optional

from logging_config import api_logger

# Lazy imports for heavy libraries — only loaded when actually needed
_docx = None
_pptx = None
_openpyxl = None
_PIL = None


def _get_docx():
    global _docx
    if _docx is None:
        import docx as _docx
    return _docx


def _get_pptx():
    global _pptx
    if _pptx is None:
        import pptx as _pptx
    return _pptx


def _get_openpyxl():
    global _openpyxl
    if _openpyxl is None:
        import openpyxl as _openpyxl
    return _openpyxl


def _get_pil():
    global _PIL
    if _PIL is None:
        from PIL import Image as _PIL
    return _PIL


# Extensions treated as plain text (read as-is)
TEXT_EXTENSIONS = {".txt", ".md", ".json", ".csv"}

# Extensions requiring binary parsing
BINARY_DOC_EXTENSIONS = {".docx", ".pptx", ".xlsx"}

# Image extensions (passed to LLM as vision content)
IMAGE_EXTENSIONS = {".jpg", ".jpeg", ".png"}

# Archive extensions
ARCHIVE_EXTENSIONS = {".zip"}

# All supported extensions
SUPPORTED_EXTENSIONS = TEXT_EXTENSIONS | BINARY_DOC_EXTENSIONS | IMAGE_EXTENSIONS | ARCHIVE_EXTENSIONS

# MIME type mapping for images
IMAGE_MIME_TYPES = {
    ".jpg": "image/jpeg",
    ".jpeg": "image/jpeg",
    ".png": "image/png",
}

# Maximum extracted text per file to prevent token overflow
MAX_TEXT_LENGTH = 100_000

# Maximum files to extract from a zip
MAX_ZIP_FILES = 20

# Maximum total extracted size from a zip (10 MB of text)
MAX_ZIP_TOTAL_SIZE = 10 * 1024 * 1024


@dataclass
class ParsedFile:
    """Result of parsing a single file."""
    name: str
    text_content: str = ""
    is_image: bool = False
    image_base64: str = ""
    image_media_type: str = ""
    parse_error: Optional[str] = None


@dataclass
class ParseResult:
    """Result of parsing all files in a request."""
    parsed_files: List[ParsedFile] = field(default_factory=list)
    image_files: List[ParsedFile] = field(default_factory=list)
    errors: List[str] = field(default_factory=list)


def get_extension(filename: str) -> str:
    """Get lowercase file extension."""
    return os.path.splitext(filename)[1].lower()


def parse_text_file(name: str, content: str) -> ParsedFile:
    """Parse a plain text file. Content is already text."""
    return ParsedFile(
        name=name,
        text_content=content[:MAX_TEXT_LENGTH],
    )


def _decode_base64_content(content: str) -> bytes:
    """Decode base64 content, stripping data URI prefix if present."""
    # Strip "data:...;base64," prefix if present
    if content.startswith("data:"):
        _, _, content = content.partition(",")
    return base64.b64decode(content)


def parse_docx(name: str, raw_bytes: bytes) -> ParsedFile:
    """Extract text from a .docx file."""
    try:
        docx = _get_docx()
        doc = docx.Document(io.BytesIO(raw_bytes))
        paragraphs = [p.text for p in doc.paragraphs if p.text.strip()]

        # Also extract text from tables
        table_rows = []
        for table in doc.tables:
            for row in table.rows:
                cells = [cell.text.strip() for cell in row.cells if cell.text.strip()]
                if cells:
                    table_rows.append(" | ".join(cells))

        text_parts = []
        if paragraphs:
            text_parts.append("\n".join(paragraphs))
        if table_rows:
            text_parts.append("\n[TABLE DATA]\n" + "\n".join(table_rows))

        text = "\n\n".join(text_parts)
        if not text.strip():
            return ParsedFile(name=name, text_content="[Empty document]")

        return ParsedFile(name=name, text_content=text[:MAX_TEXT_LENGTH])
    except Exception as e:
        api_logger.warning(f"[FILE_PARSER] Failed to parse docx {name}: {e}")
        return ParsedFile(name=name, text_content="", parse_error=f"Failed to parse docx: {e}")


def parse_pptx(name: str, raw_bytes: bytes) -> ParsedFile:
    """Extract text from a .pptx file."""
    try:
        pptx = _get_pptx()
        prs = pptx.Presentation(io.BytesIO(raw_bytes))
        slides_text = []

        for i, slide in enumerate(prs.slides, 1):
            slide_parts = [f"[Slide {i}]"]
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for paragraph in shape.text_frame.paragraphs:
                        text = paragraph.text.strip()
                        if text:
                            slide_parts.append(text)
                if shape.has_table:
                    for row in shape.table.rows:
                        cells = [cell.text.strip() for cell in row.cells if cell.text.strip()]
                        if cells:
                            slide_parts.append(" | ".join(cells))
            if len(slide_parts) > 1:  # More than just the header
                slides_text.append("\n".join(slide_parts))

        text = "\n\n".join(slides_text)
        if not text.strip():
            return ParsedFile(name=name, text_content="[Empty presentation]")

        return ParsedFile(name=name, text_content=text[:MAX_TEXT_LENGTH])
    except Exception as e:
        api_logger.warning(f"[FILE_PARSER] Failed to parse pptx {name}: {e}")
        return ParsedFile(name=name, text_content="", parse_error=f"Failed to parse pptx: {e}")


def parse_xlsx(name: str, raw_bytes: bytes) -> ParsedFile:
    """Extract data from a .xlsx file."""
    try:
        openpyxl = _get_openpyxl()
        wb = openpyxl.load_workbook(io.BytesIO(raw_bytes), read_only=True, data_only=True)
        sheets_text = []

        for sheet_name in wb.sheetnames:
            ws = wb[sheet_name]
            rows_text = [f"[Sheet: {sheet_name}]"]
            row_count = 0
            for row in ws.iter_rows(values_only=True):
                cells = [str(c) if c is not None else "" for c in row]
                # Skip fully empty rows
                if any(c for c in cells):
                    rows_text.append(" | ".join(cells))
                    row_count += 1
                    if row_count >= 2000:  # Cap rows per sheet
                        rows_text.append(f"... ({row_count}+ rows, truncated)")
                        break
            if row_count > 0:
                sheets_text.append("\n".join(rows_text))

        wb.close()

        text = "\n\n".join(sheets_text)
        if not text.strip():
            return ParsedFile(name=name, text_content="[Empty spreadsheet]")

        return ParsedFile(name=name, text_content=text[:MAX_TEXT_LENGTH])
    except Exception as e:
        api_logger.warning(f"[FILE_PARSER] Failed to parse xlsx {name}: {e}")
        return ParsedFile(name=name, text_content="", parse_error=f"Failed to parse xlsx: {e}")


def parse_image(name: str, raw_bytes: bytes) -> ParsedFile:
    """Prepare an image for LLM vision input."""
    ext = get_extension(name)
    media_type = IMAGE_MIME_TYPES.get(ext, "image/jpeg")

    try:
        # Validate the image can be opened
        pil = _get_pil()
        img = pil.open(io.BytesIO(raw_bytes))
        width, height = img.size
        img.close()

        b64 = base64.b64encode(raw_bytes).decode("utf-8")

        return ParsedFile(
            name=name,
            text_content=f"[Image: {name}, {width}x{height} pixels, {media_type}]",
            is_image=True,
            image_base64=b64,
            image_media_type=media_type,
        )
    except Exception as e:
        api_logger.warning(f"[FILE_PARSER] Failed to parse image {name}: {e}")
        return ParsedFile(name=name, text_content="", parse_error=f"Failed to parse image: {e}")


def parse_zip(name: str, raw_bytes: bytes) -> List[ParsedFile]:
    """Extract and parse files from a .zip archive."""
    results = []
    try:
        with zipfile.ZipFile(io.BytesIO(raw_bytes)) as zf:
            entries = [
                info for info in zf.infolist()
                if not info.is_dir()
                and not info.filename.startswith("__MACOSX")
                and not info.filename.startswith(".")
                and get_extension(info.filename) in (TEXT_EXTENSIONS | BINARY_DOC_EXTENSIONS | IMAGE_EXTENSIONS)
            ]

            # Sort by name and cap count
            entries.sort(key=lambda e: e.filename)
            if len(entries) > MAX_ZIP_FILES:
                api_logger.info(f"[FILE_PARSER] Zip {name} has {len(entries)} files, capping at {MAX_ZIP_FILES}")
                entries = entries[:MAX_ZIP_FILES]

            total_size = 0
            for entry in entries:
                if total_size > MAX_ZIP_TOTAL_SIZE:
                    api_logger.info(f"[FILE_PARSER] Zip {name} exceeded total size limit")
                    break

                try:
                    entry_bytes = zf.read(entry.filename)
                    total_size += len(entry_bytes)
                    entry_name = f"{name}/{entry.filename}"
                    ext = get_extension(entry.filename)

                    if ext in TEXT_EXTENSIONS:
                        text = entry_bytes.decode("utf-8", errors="replace")
                        results.append(parse_text_file(entry_name, text))
                    elif ext == ".docx":
                        results.append(parse_docx(entry_name, entry_bytes))
                    elif ext == ".pptx":
                        results.append(parse_pptx(entry_name, entry_bytes))
                    elif ext == ".xlsx":
                        results.append(parse_xlsx(entry_name, entry_bytes))
                    elif ext in IMAGE_EXTENSIONS:
                        results.append(parse_image(entry_name, entry_bytes))
                except Exception as e:
                    api_logger.warning(f"[FILE_PARSER] Failed to parse zip entry {entry.filename}: {e}")
                    results.append(ParsedFile(
                        name=f"{name}/{entry.filename}",
                        parse_error=f"Failed to extract: {e}",
                    ))

    except zipfile.BadZipFile:
        api_logger.warning(f"[FILE_PARSER] Invalid zip file: {name}")
        results.append(ParsedFile(name=name, parse_error="Invalid zip file"))
    except Exception as e:
        api_logger.warning(f"[FILE_PARSER] Failed to parse zip {name}: {e}")
        results.append(ParsedFile(name=name, parse_error=f"Failed to parse zip: {e}"))

    return results


def parse_file(name: str, content: str, encoding: str = "text") -> List[ParsedFile]:
    """
    Parse a single uploaded file.

    Args:
        name: Original filename
        content: File content (text or base64-encoded)
        encoding: 'text' for plain text content, 'base64' for base64-encoded binary

    Returns:
        List of ParsedFile objects (usually 1, but zip can produce multiple)
    """
    ext = get_extension(name)

    if ext not in SUPPORTED_EXTENSIONS:
        return [ParsedFile(name=name, parse_error=f"Unsupported file type: {ext}")]

    # Text files — content is already plain text
    if ext in TEXT_EXTENSIONS and encoding == "text":
        return [parse_text_file(name, content)]

    # For all binary formats, decode base64 first
    if encoding == "base64":
        try:
            raw_bytes = _decode_base64_content(content)
        except Exception as e:
            return [ParsedFile(name=name, parse_error=f"Failed to decode base64: {e}")]
    else:
        # Text encoding but binary extension — try to handle gracefully
        # This happens if frontend sent as text (older clients)
        if ext in TEXT_EXTENSIONS:
            return [parse_text_file(name, content)]
        # Can't parse binary from text content
        return [ParsedFile(name=name, parse_error=f"Binary file {ext} requires base64 encoding")]

    # Route to format-specific parser
    if ext == ".docx":
        return [parse_docx(name, raw_bytes)]
    elif ext == ".pptx":
        return [parse_pptx(name, raw_bytes)]
    elif ext == ".xlsx":
        return [parse_xlsx(name, raw_bytes)]
    elif ext in IMAGE_EXTENSIONS:
        return [parse_image(name, raw_bytes)]
    elif ext == ".zip":
        return parse_zip(name, raw_bytes)
    else:
        return [ParsedFile(name=name, parse_error=f"No parser for: {ext}")]


def parse_all_files(files: list) -> ParseResult:
    """
    Parse all uploaded files from a discovery request.

    Args:
        files: List of FileData-like objects with name, content, encoding, type attributes

    Returns:
        ParseResult with separated text files and image files
    """
    result = ParseResult()

    for f in files:
        encoding = getattr(f, "encoding", None) or "text"
        parsed_list = parse_file(f.name, f.content, encoding)

        for parsed in parsed_list:
            if parsed.parse_error and not parsed.text_content and not parsed.is_image:
                result.errors.append(f"{parsed.name}: {parsed.parse_error}")
                api_logger.warning(f"[FILE_PARSER] {parsed.name}: {parsed.parse_error}")
                continue

            if parsed.is_image:
                result.image_files.append(parsed)
            else:
                result.parsed_files.append(parsed)

    api_logger.info(
        f"[FILE_PARSER] Parsed {len(result.parsed_files)} text files, "
        f"{len(result.image_files)} images, {len(result.errors)} errors"
    )
    return result
